"""
SchuLyf — 15-slide academic-defense deck (python-pptx).
Assembles the StarUML-style PlantUML diagrams (uml/), the brand charts
(charts/), and real screenshots (../report/screenshots/) into a native,
editable .pptx with a consistent emerald master layout.
Output: plan/deck/build/SchuLyf-Presentation.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))
UML = os.path.join(HERE, "uml")
CH = os.path.join(HERE, "charts")
SHOT = os.path.abspath(os.path.join(HERE, "..", "report", "screenshots"))
OUT = os.path.join(HERE, "build", "SchuLyf-Presentation.pptx")
os.makedirs(os.path.dirname(OUT), exist_ok=True)

FONT = "Segoe UI"
# brand palette
EM900 = RGBColor(0x06, 0x4E, 0x3B)
EM800 = RGBColor(0x06, 0x5F, 0x46)
EM700 = RGBColor(0x04, 0x78, 0x57)
EM600 = RGBColor(0x05, 0x96, 0x69)
EM50 = RGBColor(0xEC, 0xFD, 0xF5)
EM100 = RGBColor(0xD1, 0xFA, 0xE5)
INK = RGBColor(0x0F, 0x17, 0x2A)
INK2 = RGBColor(0x47, 0x55, 0x69)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
HAIR = RGBColor(0xE2, 0xE8, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = 13.333, 7.5
BLANK = prs.slide_layouts[6]


def _run(p, text, size, color, bold=False, italic=False):
    r = p.add_run(); r.text = text
    f = r.font
    f.size = Pt(size); f.bold = bold; f.italic = italic
    f.name = FONT; f.color.rgb = color
    return r


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    return tf


def rect(slide, x, y, w, h, fill, line=None, line_w=1.0, radius=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if radius:
        try:
            shp.adjustments[0] = 0.12
        except Exception:
            pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def bg_white(slide):
    rect(slide, -0.1, -0.1, SW + 0.2, SH + 0.2, WHITE)


def base_slide(n):
    slide = prs.slides.add_slide(BLANK)
    bg_white(slide)
    # footer hairline + text
    rect(slide, 0.6, 7.02, SW - 1.2, 0.012, HAIR)
    tf = textbox(slide, 0.6, 7.08, 8, 0.3)
    _run(tf.paragraphs[0], "SchuLyf · Student Management System", 10, MUTED)
    tf2 = textbox(slide, SW - 2.6, 7.08, 2.0, 0.3)
    p = tf2.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    _run(p, f"{n:02d} / 15", 10, MUTED)
    return slide


def header(slide, kicker, title):
    tf = textbox(slide, 0.62, 0.42, 12, 0.32)
    _run(tf.paragraphs[0], kicker.upper(), 12.5, EM700, bold=True)
    tf2 = textbox(slide, 0.6, 0.74, 12.1, 0.8)
    _run(tf2.paragraphs[0], title, 29, INK, bold=True)
    rect(slide, 0.63, 1.52, 1.5, 0.06, EM600)


def place(slide, path, x, y, w, h):
    with Image.open(path) as im:
        iw, ih = im.size
    ar = iw / ih
    if ar > w / h:
        dw, dh = w, w / ar
    else:
        dh, dw = h, h * ar
    px, py = x + (w - dw) / 2, y + (h - dh) / 2
    slide.shapes.add_picture(path, Inches(px), Inches(py), Inches(dw), Inches(dh))


def caption(slide, x, y, w, text, size=13.5, color=INK2, align=PP_ALIGN.CENTER):
    tf = textbox(slide, x, y, w, 0.6)
    p = tf.paragraphs[0]; p.alignment = align
    _run(p, text, size, color)


def bullets(slide, x, y, w, h, items, size=15, gap=10):
    tf = textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.MIDDLE)
    for i, (lead, rest) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        _run(p, "●  ", size, EM600, bold=True)
        _run(p, lead, size, INK, bold=True)
        if rest:
            _run(p, rest, size, INK2)


def figslide(n, kicker, title, img, cap=None, box=(0.6, 1.75, 12.13, 4.95)):
    s = base_slide(n); header(s, kicker, title)
    x, y, w, h = box
    if cap:
        h -= 0.5
    place(s, img, x, y, w, h)
    if cap:
        caption(s, x, y + h + 0.08, w, cap)
    return s


def twofig(n, kicker, title, left, right, lcap=None, rcap=None):
    s = base_slide(n); header(s, kicker, title)
    y, h = 1.8, 4.55
    place(s, left, 0.55, y, 6.05, h)
    place(s, right, 6.75, y, 6.05, h)
    if lcap:
        caption(s, 0.55, y + h + 0.06, 6.05, lcap, size=12)
    if rcap:
        caption(s, 6.75, y + h + 0.06, 6.05, rcap, size=12)
    return s


def figbullets(n, kicker, title, img, items, imgw=7.0):
    s = base_slide(n); header(s, kicker, title)
    place(s, img, 0.55, 1.8, imgw, 4.7)
    bullets(s, 0.7 + imgw + 0.25, 1.9, SW - (0.7 + imgw + 0.25) - 0.6, 4.4, items)
    return s


# ── Slide 1 — title ───────────────────────────────────────────────────
def title_slide():
    s = prs.slides.add_slide(BLANK); bg_white(s)
    rect(s, -0.1, -0.1, SW + 0.2, 2.5, EM700)             # top band
    rect(s, -0.1, 2.4, SW + 0.2, 0.09, EM600)
    tf = textbox(s, 0.9, 0.6, 11, 1.2)
    _run(tf.paragraphs[0], "SchuLyf", 54, WHITE, bold=True)
    tf2 = textbox(s, 0.95, 1.72, 11, 0.5)
    _run(tf2.paragraphs[0], "A Student Management System for the university", 18, EM100)
    # project title
    tf3 = textbox(s, 0.9, 3.1, 11.5, 1.4)
    _run(tf3.paragraphs[0],
         "Digitising admissions, tamper-proof payments, exam-gating,\n"
         "course management & verifiable academic transcripts", 22, INK, bold=True)
    # identity placeholders
    rows = [("Presented by", "«Your Name»"),
            ("Institution", "«University / Faculty»"),
            ("Department / Course", "«Department · Course code»"),
            ("Supervisor", "«Supervisor name»"),
            ("Date", "«Month Year»")]
    y = 4.7
    for label, val in rows:
        tfa = textbox(s, 0.95, y, 3.0, 0.35)
        _run(tfa.paragraphs[0], label.upper(), 11, EM700, bold=True)
        tfb = textbox(s, 4.0, y, 8.2, 0.35)
        _run(tfb.paragraphs[0], val, 14, INK)
        y += 0.46
    tf4 = textbox(s, 0.9, 7.05, 11, 0.3)
    _run(tf4.paragraphs[0], "Laravel 13 · Inertia · Vue 3 · PrimeVue · MySQL · Pest", 11, MUTED)
    return s


# ── stat tiles (slide 14) ─────────────────────────────────────────────
def stat_tiles(slide, x, y, w, h, tiles):
    n = len(tiles); gap = 0.25
    tw = (w - (n - 1) * gap) / n
    for i, (num, lab) in enumerate(tiles):
        tx = x + i * (tw + gap)
        rect(slide, tx, y, tw, h, EM50, EM600, 1.2, radius=True)
        tf = textbox(slide, tx, y + 0.18, tw, h - 0.6, anchor=MSO_ANCHOR.MIDDLE)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        _run(p, num, 30, EM800, bold=True)
        tf2 = textbox(slide, tx, y + h - 0.5, tw, 0.4)
        p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
        _run(p2, lab, 12.5, INK2)


def results_slide():
    s = base_slide(14); header(s, "Results", "Implementation & results")
    stat_tiles(s, 0.6, 1.75, 12.13, 1.35,
               [("127", "routes"), ("34", "migrations"), ("28", "ADRs"),
                ("707", "tests"), ("7", "modules"), ("6", "roles")])
    place(s, os.path.join(CH, "test-growth.png"), 0.55, 3.35, 6.5, 3.35)
    # two real screenshots (right)
    place(s, os.path.join(SHOT, "03-sao-review-queue.png"), 7.35, 3.3, 5.4, 1.7)
    place(s, os.path.join(SHOT, "08-admin-dashboard.png"), 7.35, 5.05, 5.4, 1.7)
    caption(s, 7.35, 6.72, 5.4, "Real screens — SAO review queue · Admin dashboard", size=11)
    return s


def conclusion_slide():
    s = base_slide(15); header(s, "Testing · Conclusion", "Testing, a hard-won catch & what's next")
    place(s, os.path.join(CH, "ci-testpyramid.png"), 0.5, 1.75, 6.6, 4.9)
    x = 7.4; w = SW - x - 0.6
    tf = textbox(s, x, 1.85, w, 4.8)
    p = tf.paragraphs[0]
    _run(p, "Engineering rigor", 16, EM800, bold=True)
    for t in ["Test-driven; 707 automated Pest tests",
              "4 required CI checks green before every merge",
              "Per-feature quality gate (Pint · vue-tsc · ESLint · build)"]:
        pp = tf.add_paragraph(); pp.space_after = Pt(4)
        _run(pp, "●  ", 13, EM600, bold=True); _run(pp, t, 13, INK2)
    pp = tf.add_paragraph(); pp.space_before = Pt(10)
    _run(pp, "A notable catch", 16, EM800, bold=True)
    pp = tf.add_paragraph()
    _run(pp, "The final whole-branch review found a MySQL JSON "
             "key-order bug that made every genuine transcript read "
             "“invalid” — invisible to the SQLite test suite. "
             "Fixed by canonicalising the digest.", 12.5, INK2)
    pp = tf.add_paragraph(); pp.space_before = Pt(10)
    _run(pp, "Future work", 16, EM800, bold=True)
    pp = tf.add_paragraph()
    _run(pp, "SMS notification channel · analytics dashboards · "
             "mobile app · timetable & scheduling.", 12.5, INK2)
    # closing band
    rect(s, 7.4, 6.35, w, 0.5, EM700, radius=True)
    tfc = textbox(s, 7.4, 6.4, w, 0.4, anchor=MSO_ANCHOR.MIDDLE)
    pc = tfc.paragraphs[0]; pc.alignment = PP_ALIGN.CENTER
    _run(pc, "Thank you — questions welcome", 14, WHITE, bold=True)
    return s


# ── build ─────────────────────────────────────────────────────────────
title_slide()
figslide(2, "Context", "The problem", os.path.join(CH, "pain-points.png"),
         "Five recurring failures in day-to-day student administration.")
figslide(3, "Goals", "Objectives & scope", os.path.join(CH, "objectives-modules.png"),
         "Each real-world pain maps to a dedicated SchuLyf module.")
figbullets(4, "Approach", "Methodology & delivery", os.path.join(CH, "phase-timeline.png"),
           [("Laravel 13 + Inertia/Vue 3", " full-stack, session-authenticated"),
            ("Test-driven (Pest)", " every change proven by a test"),
            ("Per-feature quality gate", " Pint · types · lint · build"),
            ("Audit-driven hardening", " 34 findings remediated"),
            ("Phased, incremental delivery", "")], imgw=7.4)
figslide(5, "Design", "System architecture",
         os.path.join(UML, "component-architecture.png"),
         "Thin controllers → actions/services → Eloquent → MySQL; queued mail; Fortify auth.")
figslide(6, "Design", "Actors & use cases", os.path.join(UML, "usecase.png"),
         "Six roles plus an unauthenticated public verifier.")
figslide(7, "Design", "Domain model (UML class diagram)",
         os.path.join(UML, "class-domain.png"),
         "Core classes, attributes, operations and relationships.")
figslide(8, "Design", "Data model (crow's-foot ER)",
         os.path.join(UML, "er-datamodel.png"),
         "34 migrations; the core entities and their cardinalities.")
twofig(9, "Flows", "Admissions — decision & lifecycle",
       os.path.join(UML, "seq-decision.png"), os.path.join(UML, "state-application.png"),
       "SAO decision sequence (admit mints a student)",
       "Application state machine (born Submitted)")
twofig(10, "Flows", "Tamper-proof payments & receipts",
       os.path.join(UML, "seq-payment.png"), os.path.join(CH, "hmac-signature.png"),
       "Validate → mint an immutable, signed receipt",
       "HMAC binds identity; verify re-derives it")
twofig(11, "Security", "Verification & the security model",
       os.path.join(UML, "seq-verify.png"), os.path.join(CH, "security-layers.png"),
       "Public verify — authentic or “invalid”, no oracle",
       "Defense in depth across five layers")
twofig(12, "Flows", "Exam gating — standing & deferrals",
       os.path.join(UML, "activity-standing.png"), os.path.join(CH, "standing-thresholds.png"),
       "Standing decides exam-hall access each deadline",
       "Below threshold → at risk unless deferred")
figbullets(13, "Feature", "Academic transcripts (#71)",
           os.path.join(CH, "transcript-gpa.png"),
           [("4.0 GPA / CGPA", " credit-weighted, per semester & cumulative"),
            ("Immutable snapshot", " signed at issue; results may change later"),
            ("mpdf PDF + QR", " branded, downloadable transcript"),
            ("Public HMAC verify", " authentic-or-invalid, no oracle")], imgw=6.7)
results_slide()
conclusion_slide()

prs.save(OUT)
print("saved", OUT, "-", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
